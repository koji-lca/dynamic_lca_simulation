#!/usr/bin/env python3
"""
Generate SuMPO-branded 20-min presentation PPTX for overseas dynamic LCA trends.
Uses the official SuMPO template (sumpo-ppt.skill / shared/templates/sumpo_template.pptx).

Layout indices (per SKILL.md):
  0  1_タイトル スライド         → 表紙
  1  1_見出し 水色              → セクション区切り（水色）
  2  1_タイトルとコンテンツ 水色 → コンテンツ（水色）
  7  4_見出し 紺色              → セクション区切り（紺色）
  8  4_タイトルとコンテンツ 紺色 → コンテンツ（紺色・デフォルト）
  10 2_最終ページ 機構名         → エンドスライド（青背景）
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ── Paths ──────────────────────────────────────────────────────────────
HERE = Path(__file__).parent
TEMPLATE = HERE.parent.parent.parent / "shared" / "templates" / "sumpo_template.pptx"
OUTPUT = HERE / "presentation_海外動的LCA動向_20min.pptx"

# ── SuMPO brand colors ─────────────────────────────────────────────────
C_TITLE  = RGBColor(0x12, 0x54, 0x86)   # Chathams Blue  (タイトル強調)
C_DARK   = RGBColor(0x49, 0x4D, 0x51)   # Abbey          (テーブルヘッダー)
C_SILVER = RGBColor(0xC5, 0xC9, 0xCD)   # Silver Sand    (テーブルセル)
C_NAVY   = RGBColor(0x00, 0x57, 0xBA)   # Absolute Zero  (紺アクセント)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ── Layout indices ─────────────────────────────────────────────────────
L_COVER   = 0    # 表紙
L_SEC_BLU = 1    # セクション区切り（水色）
L_CON_BLU = 2    # コンテンツ（水色）
L_SEC_NAV = 7    # セクション区切り（紺色）
L_CON_NAV = 8    # コンテンツ（紺色） ← デフォルト
L_END_BLU = 10   # 最終ページ（青背景）


# ── Helpers ────────────────────────────────────────────────────────────

def _delete_slide(prs: Presentation, idx: int) -> None:
    """Remove slide at 0-based index from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[idx]
    rId = sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _add_slide(prs: Presentation, layout_idx: int):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _set_title(slide, text: str, size: int = 22) -> None:
    sh = slide.shapes.title
    if sh is None:
        return
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = C_TITLE


def _fill_body(slide, items: list, size: int = 16, ph_idx: int = 1) -> None:
    """Fill the body placeholder.

    items: list of str  →  level=0
           list of (str, int)  →  (text, indent_level)
    Empty string '' inserts a blank spacer paragraph.
    """
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in items:
        text  = item[0] if isinstance(item, tuple) else item
        level = item[1] if isinstance(item, tuple) else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text  = text
        p.level = level
        p.font.size = Pt(size)


def _set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _add_table(slide, headers: list, rows: list,
               x: int, y: int, width: int, height: int,
               header_size: int = 12, body_size: int = 11) -> None:
    """Add a SuMPO-styled table to the slide.

    headers: list of str column headers
    rows: list of list of str
    Colors: header=#494D51 white / alt rows=#C5C9CD
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, width, height).table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    for i, row_data in enumerate(rows):
        use_alt = (i % 2 == 0)
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            if use_alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_SILVER
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(body_size)


# ── Slide builders ─────────────────────────────────────────────────────

def s01_cover(prs: Presentation):
    slide = _add_slide(prs, L_COVER)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "海外における動的LCAの動きとその意味"
        run.font.size = Pt(28)
        run.font.bold = True
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            lines = [
                "～なぜ導入され、何が変わったのか～",
                "",
                "林野庁委託「建築物LCA制度化に向けた調査」",
                "政策・制度調査（第1係）",
                "2026年5月（ドラフト）｜想定20分・14枚",
            ]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
    _set_notes(slide, "フランスを中心に米国・EUの文脈も含め「海外で何が起きているか」を整理します。")
    return slide


def s02_goals(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "本日のゴール")
    _fill_body(slide, [
        "【本日お伝えする3点】",
        ("① なぜ — 海外（特にフランス）で動的LCAが政策の中心に入ったのか", 0),
        ("② どう — 制度化されたか（RE2020の要点）", 0),
        ("③ 導入後 — 産業・設計・数値にどんな影響が出ているか", 0),
        "",
        "【本日扱わないこと】",
        ("算定式の詳細証明 → 第2係", 1),
        ("INIES・DB構造の詳細 → 第3係", 1),
    ])
    # Key Message as a visually distinct text box
    txBox = slide.shapes.add_textbox(
        Emu(400000), Emu(5500000), Emu(9900000), Emu(1300000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "★ Key Message"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "「フランスは建築LCAに『時間』を制度化した世界で唯一の国 — その意味と日本への示唆を整理する」"
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = C_NAVY
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = C_SILVER
    _set_notes(slide, "聴き取りの軸を3点に絞ります。算定式詳細・INIES構造は他係が担当。")
    return slide


def s03_global_trends(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "世界の流れ：建築の炭素を「見える化」する時代")
    _fill_body(slide, [
        "共通する政策課題",
        ("運用エネルギー効率化: 既に各国で進展（断熱・設備・再エネ）", 1),
        ("エンボディド炭素: 建設・部材・解体が残る大きな削減余地", 1),
        ("ネットゼロ整合: 建物はライフサイクル全体で説明責任を求められる", 1),
        "",
        "各国の制度化の「入り口」（論点は共通）",
        ("フランス: RE2020で建築物LCA（ACV）を義務化、簡易動的LCAを採用（2022〜）", 1),
        ("EU: EPBD改正で2030年以降新築にWLC開示義務化。ISO21391-1:2025が2025年成立", 1),
        ("米国: Federal Buy Clean・LEED v5でエンボディド炭素LCAを前提化（2025〜）", 1),
        "",
        "→ 「建築の炭素」は材料・時間・データの問題として制度化が始まっている",
    ])
    _set_notes(slide, "入り口は違うが、材料・時間・データが共通課題です（2分）。")
    return slide


def s04_static_limits(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "静的LCAだけでは足りない理由")
    _fill_body(slide, [
        "従来のGWP100（静的LCA）",
        ("ライフサイクル全体の排出を一度に足し合わせ、100年GWPで換算", 1),
        ("国際標準（EN15804・ISO14067等）で比較しやすいが、時間構造を捨てている", 1),
        "",
        "建築・木材で生じる「時間の論点」",
        ("一時的炭素貯蔵（木材固定）: 静的LCAは考慮しない", 1),
        ("  → 「CLTはコンクリートと正味ゼロ」という物理的に誤った答えが出る", 2),
        ("遅延排出（解体・焼却が数十年後）: 静的では今日放出と同一視", 1),
        ("  → 将来排出の気候影響を42〜53%過大評価している", 2),
        ("生物由来CO₂: 吸収・放出のタイミングが手法により異なる", 1),
        "",
        "学術的出発点: Levasseur et al. (2010/2013) — 排出の発生時刻を保持し累積放射強制力で評価",
        "→ フランスはこれを規制向けに「簡略化」して採用した",
    ])
    _set_notes(slide, "静的は比較可能性の強み。動的は時間の論点に応える試みです（4分）。")
    return slide


def s05_global_map(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "海外マップ：どこまで「動的」か")
    _fill_body(slide, [
        "規制で動的LCAを採用",
        ("フランス RE2020（2022〜）: 簡易動的LCA・義務", 1),
        "",
        "LCA枠組み・任意指標",
        ("EU Level(s) / EN15978: 建築LCAの共通言語；動的は国別実装", 1),
        ("ISO 21391-1:2025: 動的バイオジェニック炭素評価の初国際規格", 1),
        "",
        "調達・認証でLCA強化",
        ("米国 Buy Clean / LEED v5 / EPA低炭素建材ラベル（2024〜）", 1),
        "",
        "研究・ツール",
        ("カナダ CIRAIG DynCO₂（RE2020係数の理論的源流）など", 1),
        "",
        "重要: 「学術の完全版動的LCA」と「規制の簡易版」を意図的に分けて理解すること",
    ])
    _set_notes(slide, "動的LCA＝学術概念と規制の簡易版を分けて見ることが重要です（6分）。")
    return slide


def s06_re2020(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "フランス RE2020：何が変わったか")
    _fill_body(slide, [
        "RE2020の位置づけ（Cerema Guide 2024）",
        ("RT2012（熱・エネルギー中心）→ RE2020（エネルギー＋環境性能＝第3の柱）", 1),
        ("建築物のライフサイクル（ステージ1〜5＋輸送）で環境影響を評価", 1),
        ("規制対象の気候指標: GES（温室効果ガス）— kgCO₂eq/m²", 1),
        "",
        "規制値と段階適用スケジュール（住宅・標準規模の代表値）",
        ("Phase 1（2022）: 戸建 640 / 集合 840 / 事務所 740 kgCO₂eq/m²", 1),
        ("Phase 2（2025）: 各 −15%", 1),
        ("Phase 3（2028）: 各 −25%（DCF係数見直し予定）", 1),
        ("Phase 4（2031）: 各 −30〜40%", 1),
        "",
        "データ基盤: INIES（建材環境DB）、FDES（フランスのEPD）",
        "※Arrêté 2021 Annexe II 数値は原文照合中",
    ])
    _set_notes(slide, "制度の骨格と段階適用スケジュールを整理します（7.5分）。")
    return slide


def s07a_physics(prs: Presentation):
    """Slide 7a: Physical basis + policy consequence."""
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "なぜ動的LCAか①：物理的根拠と政策的帰結")
    _fill_body(slide, [
        "【物理的根拠】なぜ50年後の排出は「軽い」のか",
        ("今日放出の1 kg CO₂: 100年間気候に影響し続ける → DCF(0) = 1.00", 1),
        ("50年後放出の1 kg CO₂: 残り50年のみ影響 → DCF(50yr) = 0.58（−42%）", 1),
        ("静的LCAの前提「今日の1 kg ≡ 50年後の1 kg」は物理的に誤り", 1),
        ("  根拠: IPCC AR4 Bern2.5CCモデル（Benoist 2009・Solinnen 2018）", 2),
        "",
        "【政策的帰結】動的LCAで木材の評価が逆転する",
        ("CLT（木材）: 静的 +98 kgCO₂eq/m³（排出）→ 動的 −171（吸収）★+269の評価差", 1),
        ("RC造（コンクリート）: 静的 +300 → 動的 +295（炭素貯蔵がないためほぼ不変）", 1),
        ("※数値は概算（第1編報告書§4.7）。INIES/FDES一次データ照合中", 1),
        "",
        "→ 設計者が木材を選ぶ合理的根拠を「数字」で初めて示せるツールが誕生",
    ])
    _set_notes(slide, "物理的根拠（9.5分）→政策的帰結（10.5分）。CLT数値は概算注記を確認してください。")
    return slide


def s07b_necessity(prs: Presentation):
    """Slide 7b: 4 convergent necessities + institutional design."""
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "なぜ動的LCAか②：4つの必然性と制度設計の選択")
    _fill_body(slide, [
        "【4つの収束した必然性】なぜフランスで・なぜ今だったか",
        ("① 科学的必要性: Benoist(2009)→Solinnen(2018) で係数確立。物理的根拠がある", 1),
        ("② 政策的整合: SNBC 2050目標は建材炭素規制なしに達成不可。静的LCAでは誤ったシグナル", 1),
        ("③ 産業戦略: 国土31%が森林のフランスに「評価ツールの失敗が産業を不利にしていた」状況の解消", 1),
        ("④ 国際先行者戦略: ISO 21391-1:2025（AFNOR主導）・EPBD 2024への影響力確保", 1),
        "",
        "【制度設計の実務的選択】なぜ「簡易」線形近似か",
        ("透明性: 係数表で誰でも再現可能（Excelで計算できる）", 1),
        ("一貫性: パラメータ固定で規制の予測可能性が高い", 1),
        ("移行可能性: 既存LCAソフト＋係数表で実装。開発コスト低", 1),
        ("線形近似式: 1 − 0.00842·t（0〜50年）; 0.58固定（50年超）", 1),
        ("2028年Phase 3での非線形DCF（ISO 21391-1方式）移行を計画的に留保したアダプティブ設計", 1),
    ])
    _set_notes(slide, "4つの必然性（11分）→制度設計の選択（11.5分）。")
    return slide


def s08_dcf_numbers(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "数値で見る「時間地平」の威力")
    # Minimal intro: one line only to leave room for table
    _fill_body(slide, [
        "制度が選ぶ観測年数＝政策メッセージ。出典: Ventura & Feraille (2021); RE2020 Dynamic LCA資料",
    ])
    # DCF comparison table — starts high enough to avoid body overlap
    headers = ["排出時点 t", "RE2020線形近似", "真のDCF (DO=100yr)", "Ventura提案", "RE2020過大誤差"]
    rows = [
        ["t=0（建設時）",    "1.000", "1.000", "1.000", "0%"],
        ["t=10年",           "0.916", "0.915", "0.951", "+3.9%"],
        ["t=25年",           "0.790", "0.792", "0.858", "+8.3%"],
        ["t=50年（設計寿命）", "0.579", "0.606", "0.720", "+24%  ★"],
        ["t>50年（固定）",   "0.580", "0.427 (t=75yr)", "0.617 (t=75yr)", "最大53% ★★"],
    ]
    _add_table(
        slide, headers, rows,
        x=Emu(400000), y=Emu(2200000),
        width=Emu(9900000), height=Emu(3200000),
    )
    # Footnote textbox below table
    note_box = slide.shapes.add_textbox(
        Emu(400000), Emu(5600000), Emu(9900000), Emu(600000)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("★ t=50年: RE2020係数は物理的DCFより最大24%過大  "
                "★★ t>50年: RE2020固定値0.580 vs 真DCF 0.427 — 差53%（2028年Phase 3でISO21391-1方式に移行予定）")
    run.font.size = Pt(10)
    run.font.color.rgb = C_TITLE
    _set_notes(slide, "13分。★=係数過大問題の核心。2028年Phase 3で非線形DCF（ISO21391-1）へ移行予定。")
    return slide


def s09_impacts(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "RE2020導入後の影響")
    _fill_body(slide, [
        "1. 設計・資材選択",
        ("CLT: 静的+98 → 動的−171 kgCO₂eq/m²（評価が逆転）", 1),
        ("設計者が木材を選ぶ合理的根拠を「数字」で初めて示せるツールが誕生", 1),
        ("静的LCAでは見えなかった解体時排出の遅れが、木材に相対的に有利に機能", 1),
        "",
        "2. 産業・データ",
        ("FDES/EPD・INIESの整備需要が増加", 1),
        ("セメント・金属産業: EN15804/PEF/ISO14067（一時貯蔵除外）との不整合を主張し反対", 1),
        ("2025年7月 Rivaton報告書: ACVルール・データ品質・コストを評価、12項目勧告", 1),
        "",
        "3. 政策目標との接続",
        ("2025/2028/2031の段階削減を動的GES指標でモニタリング", 1),
        ("動的LCA採用＝木造義務化ではない（設計インセンティブとして機能）", 1),
        "",
        "4. 国際標準との緊張（継続中）",
        ("フランス独自簡易動的は国内比較に有効、国際EPDとの横並びに説明が必要", 1),
    ])
    _set_notes(slide, "設計・産業・政策・国際標準の4軸で整理します（15分）。")
    return slide


def s10_usa(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "アメリカ：同じ課題、異なる答え")
    _fill_body(slide, [
        "米国は「動的LCAを建築法で義務化」していない（仮説）",
        ("連邦建築法の分権: 全国一律算定より調達・認証で寄せる方式", 1),
        ("比較可能性・訴訟リスク: 建材産業が手法の公平性を重視（仏と同型の対立構造）", 1),
        ("研究先行: 動的・バイオジェニックは学術・ツール段階; 規制は静的GWP+EPDから", 1),
        "",
        "エンボディド炭素を「国家・市場で押し上げる」動きは急速",
        ("Federal Buy Clean: 連邦調達で低炭素建材・EPDを優先（材料GWP開示が中心）", 1),
        ("EPA低炭素建材ラベル（IRA 60116）: PCR・バイオジェニック開示要件（2024）", 1),
        ("LEED v5（2025〜）: エンボディド炭素LCAを前提化 / 州・都市条例は分断的", 1),
        "",
        "それでも「フランスと同じ論点」が来ている",
        ("大規模木造・マスティンバー推進 ↔ バイオジェニック炭素をどう数えるか", 1),
        ("認証間で貯蔵・EoLの扱いがずれる問題（LEED/BREEAM比較研究）", 1),
    ])
    _set_notes(slide, "米国は調達・認証で押し上げ。動的統一義務は未整備（17分）。")
    return slide


def s11_comparison(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "フランス vs 米国：同じ論点・異なる制度化")
    # Comparison table
    headers = ["観点", "フランス（RE2020）", "米国（Buy Clean / LEED v5）"]
    rows = [
        ["義務の形",      "建築環境規制（新築・法定）",         "調達・認証・州法の組み合わせ"],
        ["LCAの位置",     "GES指標で法定",                     "LEED v5で前提化；建築コード未統一"],
        ["時間軸",        "簡易動的LCA（FRE2020）を既定手法",   "原則静的GWP；動的は個別研究段階"],
        ["木材・炭素貯蔵", "制度上クレジット化（StockC計上）",  "開示・事例依存；統一クレジットなし"],
        ["産業政治",      "セメント vs 木材の公開論戦",         "Buy Cleanで鋼・コンクリートが主戦場"],
        ["データ",        "INIES / FDES（フランスEPD）",        "EPA PCR / 民間LCAツール"],
        ["日本への示唆",  "年数・起点を制度で先に決める",       "義務化段階とDB整備を分けて設計"],
    ]
    _add_table(
        slide, headers, rows,
        x=Emu(400000), y=Emu(1800000),
        width=Emu(9900000), height=Emu(5200000),
    )
    _set_notes(slide, "18分。比較表で「同じ論点・違う制度化の入り口」を示します。")
    return slide


def s12_eu(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "EU・その他（補足）")
    _fill_body(slide, [
        "EU Level(s) / EN15978",
        ("建築のサステナビリティ指標の共通フレーム", 1),
        ("動的LCAの統一義務はフランスが先行；EUはWLC開示義務化に向けて整備中", 1),
        "",
        "ISO 21391-1:2025（2025年成立）",
        ("建材の動的バイオジェニック炭素評価に関する初の国際規格", 1),
        ("非線形DCFを採用し、RE2020線形近似より精度を改善", 1),
        ("フランス（AFNOR）が主導、EPBD2024への影響力を確保", 1),
        "",
        "研究・ツール（カナダ CIRAIG 等）",
        ("DynCO₂: 動的炭素フットプリント計算機 — RE2020係数の理論的源流", 1),
        ("日本が参照するなら「学術の完全版」と「規制の簡易版」を意図的に分離すること", 1),
    ])
    _set_notes(slide, "補足1分。詳細はNo.3（Level(s)）で深掘り予定。")
    return slide


def s13b_roadmap(prs: Presentation):
    """Slide 13b: Japan implementation roadmap as table."""
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "日本の動的LCA制度化：中長期ロードマップ")
    _fill_body(slide, [
        "フランスの教訓: E+C-実証（2016〜2021, 2,200棟）→Solinnen報告（2018）→RE2020施行（2022）— 6年プロセス",
    ])
    # Roadmap as table
    headers = ["ステップ", "時期", "主なアクション", "前提条件（必要）"]
    rows = [
        ["Step 1", "〜2027",      "ISO 21391-1:2025 JIS化着手\nIDEA時系列フロー追加設計",   "時間地平（DO）の政策的合意"],
        ["Step 2", "2027〜2028", "実証プログラム（300〜500棟）\n適用対象・時間地平の合意形成",  "建材DB整備（IDEA拡張）"],
        ["Step 3", "2028〜2029", "パラメータ実検\nDO・起点・統一デフォルト値",               "算定エンジンの国産化"],
        ["Step 4", "2029〜2030", "Ic上限設定（段階値）\n林野庁・国土交通省と制度設計調整", "実証データの裏付け"],
        ["Step 5", "2030〜2031", "段階適用開始\n建物種別・規模別フェーズイン",          "国際整合（EPBD 2030）確認"],
    ]
    _add_table(
        slide, headers, rows,
        x=Emu(400000), y=Emu(2200000),
        width=Emu(9900000), height=Emu(4500000),
        body_size=12,
    )
    _set_notes(slide, "19〜19.5分。フランスの6年プロセスを日本版5ステップに対応付ける。")
    return slide


def s13_japan(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "日本の建築物LCA制度化への示唆（優先課題）")
    _fill_body(slide, [
        "制度設計のための最優先対応3点",
        ("① 時間地平の合意形成", 0),
        ("  DO=100年 vs 50年で評価差が2倍近く変わる", 1),
        ("  → 先に定めなければ規制値議論自体が無意味", 1),
        ("② ISO 21391-1:2025のJIS化", 0),
        ("  JIS化に3〜5年要する。2030年EPBD義務化に国際整合できなくなるリスク", 1),
        ("③ 建材LCAデータベース整備", 0),
        ("  IDEAへの時系列フロー追加が最低限要件（現状では動的LCAは机上論）", 1),
        "",
        "海外から学べる3点",
        ("時間の設計が本体: 起点・DOの制度内確定が期待値を決定（フランスの事後論争が証明）", 1),
        ("段階導入の現実解: 静的＋簡易動的ハイブリッドで10年実証後に規制化（仏方式）", 1),
        ("林野・木材政策との接続: 評価手法が木造インセンティブになり得る。持続可能性証明とセット設計", 1),
    ])
    _set_notes(slide, "日本は両国の「中間」も可能。まず時間地平の合意を推奨します（19分）。")
    return slide


def s14_summary(prs: Presentation):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "まとめ・今後の調査方針")
    _fill_body(slide, [
        "唯一のテイクアウェイ",
        ("「日本が動的LCA制度化を目指すなら、規制値議論の前に:", 1),
        ("①時間地平の合意 → ②ISO21391-1のJIS化 → ③建材DB整備 の3つを先に達成すること。", 1),
        ("フランスはこの3つを整えてから規制化した（E+C-実証2016〜2021→RE2020施行2022）。」", 1),
        "",
        "3行まとめ",
        ("海外 = 建築のエンボディド炭素を、LCAで制度・市場に載せるフェーズに入った", 1),
        ("フランス = 唯一、簡易動的LCAを建築規制の中心に置いた先進事例", 1),
        ("日本への示唆 = ①時間地平合意 → ②ISO JIS化 → ③DB整備 の順に先行投資", 1),
        "",
        "第1係の次ステップ",
        ("Arrêté・Annexe 条文照合（No.1残件）→ 第1係内完結", 1),
        ("制度導入背景・林業戦略（No.4〜7）→ 次期成果物", 1),
        ("INIES/FICES構造分析要件の整理 → 第3係への引き継ぎ", 1),
    ])
    _set_notes(slide, "質疑: 日本導入は時間地平とデータ基盤の合意が先（20分）。")
    return slide


def s15_end(prs: Presentation):
    return _add_slide(prs, L_END_BLU)


# ── Main ───────────────────────────────────────────────────────────────

def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(
            f"SuMPO template not found: {TEMPLATE}\n"
            "Extract sumpo-ppt.skill and copy assets/template.pptx to shared/templates/sumpo_template.pptx"
        )

    prs = Presentation(str(TEMPLATE))

    # Remove all demo slides from the template (reverse order to keep indices stable)
    n_existing = len(prs.slides)
    for i in range(n_existing - 1, -1, -1):
        _delete_slide(prs, i)

    builders = [
        s01_cover,
        s02_goals,
        s03_global_trends,
        s04_static_limits,
        s05_global_map,
        s06_re2020,
        s07a_physics,
        s07b_necessity,
        s08_dcf_numbers,
        s09_impacts,
        s10_usa,
        s11_comparison,
        s12_eu,
        s13_japan,
        s13b_roadmap,
        s14_summary,
        s15_end,
    ]

    for fn in builders:
        fn(prs)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
