"""build_survey1_ppt.py — Survey 1 SuMPO PPTX builder (SKILL.md compliant)"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE = '/tmp/sumpo_skill/extracted/sumpo-ppt-v2/assets/template.pptx'
OUTDIR = os.path.join(os.path.dirname(__file__), 'slides')
OUTPUT  = os.path.join(OUTDIR, 'survey1_政策制度調査_SuMPO.pptx')

# SuMPO colors (SKILL.md)
CHATHAMS = RGBColor(0x12,0x54,0x86)
ABSOLUTE = RGBColor(0x00,0x57,0xBA)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
BLACK    = RGBColor(0x00,0x00,0x00)
DGRAY    = RGBColor(0x44,0x44,0x44)

# Layout indices (SKILL.md)
L_COVER, L_SEC, L_CON, L_END = 0, 7, 8, 10

# Content area EMU coords (from template Layout 8)
BL, BT, BW, BH = 510948, 1381566, 9721187, 5278273

def clear_slides(prs):
    R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    lst = prs.slides._sldIdLst
    for s in list(lst):
        rid = s.get(f'{{{R}}}id')
        if rid:
            try: prs.part.drop_rel(rid)
            except: pass
        lst.remove(s)

def set_title(slide, text, sz=20, color=CHATHAMS, bold=True):
    sh = slide.shapes.title
    if not sh: return
    sh.text_frame.clear()
    p = sh.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color

def body(slide, lines, l=BL, t=BT, w=BW, h=BH, sz=13):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        bold = line.startswith('■') or line.startswith('【') or line.startswith('▶')
        color = CHATHAMS if line.startswith('■') else (ABSOLUTE if line.startswith('▶') or line.startswith('【') else (DGRAY if line.startswith('  ') else BLACK))
        r = p.add_run(); r.text = line.lstrip()
        r.font.size = Pt(sz - (1 if line.startswith('  ') else 0))
        r.font.bold = bold; r.font.color.rgb = color
        p.space_after = Pt(1)

def cover(prs, title, sub, date):
    sl = prs.slides.add_slide(prs.slide_layouts[L_COVER])
    set_title(sl, title, sz=24, color=CHATHAMS)
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]; r = p.add_run()
            r.text = sub + '\n' + date; r.font.size = Pt(14); r.font.color.rgb = DGRAY

def section(prs, num, title, sub=''):
    sl = prs.slides.add_slide(prs.slide_layouts[L_SEC])
    set_title(sl, f'{num}　{title}', sz=26, color=WHITE)
    if sub:
        for ph in sl.placeholders:
            if ph.placeholder_format.idx not in (0, 8):
                ph.text = sub
                for p in ph.text_frame.paragraphs:
                    for r in p.runs: r.font.size = Pt(15); r.font.color.rgb = WHITE

def content(prs, title, lines, sz=13):
    sl = prs.slides.add_slide(prs.slide_layouts[L_CON])
    set_title(sl, title); body(sl, lines, sz=sz)

def two_col(prs, title, lt, ll, rt, rl):
    HW = (BW - 200000) // 2
    sl = prs.slides.add_slide(prs.slide_layouts[L_CON])
    set_title(sl, title)
    # left header + body
    tb = sl.shapes.add_textbox(Emu(BL), Emu(BT), Emu(HW), Emu(280000))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = lt; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = CHATHAMS
    body(sl, ll, l=BL, t=BT+300000, w=HW, h=BH-300000, sz=12)
    # right header + body
    RL = BL + HW + 200000
    tb2 = sl.shapes.add_textbox(Emu(RL), Emu(BT), Emu(HW), Emu(280000))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = rt; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = CHATHAMS
    body(sl, rl, l=RL, t=BT+300000, w=HW, h=BH-300000, sz=12)

def end(prs):
    prs.slides.add_slide(prs.slide_layouts[L_END])

# ── Slide content ──────────────────────────────────────────────────────────────
def build():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    cover(prs,
        '調査ブロック①　政策及び制度調査報告',
        'フランスRE2020における動的LCA制度の全体構造と導入背景',
        '一般社団法人サステナブル経営推進機構（SuMPO）　2026年6月')

    content(prs, '目次', [
        '① 建築環境性能評価制度の全体構造         No.1-3・8',
        '  • RE2020の制度的位置付けと3評価指標',
        '  • 段階的seuil値・Level(s)・EN15978との関係',
        '② ダイナミックLCAの要件整理              No.9-11',
        '  • 動的LCA制度要件・DCF・FRE2020(t)算定根拠',
        '  • 段階施行スケジュールと主要パラメータ',
        '③ 生物由来炭素・炭素貯蔵の制度的位置付け  No.12-14',
        '  • StockC算定方法・評価期間（DO）の根拠',
        '  • 解体・再利用シナリオ',
        '④ 制度導入の背景および検討プロセス        No.4-7',
        '  • 4条件収束・15年の経緯・Loi Climat',
        '  • 産業界の対立と調停・代替案比較・国際戦略',
        '⑤ 日本への示唆・まとめ・参考文献',
    ])

    # ── Section 1 ──────────────────────────────────────────────────────────────
    section(prs, '①', '建築環境性能評価制度の全体構造', 'No.1（ア）/ No.2（イ）/ No.3（ウ）/ No.8（エ）')

    content(prs, 'RE2020の制度的位置付けと3評価指標', [
        '■ RE2020 = 世界初の建物LCA義務化規制（フランス、2022年施行）',
        '  Réglementation Environnementale 2020。新築建物を対象。',
        '',
        '■ 3つの評価指標（三本の柱）',
        '【Ic_construction】エンボディド炭素 [kgCO₂eq/m²]',
        '  A〜C モジュールGHG + StockC（バイオジェニック炭素クレジット）',
        '  → 動的LCA係数 FRE2020(t) を適用',
        '【Ic_énergie】運用エネルギー由来CO₂ [kgCO₂eq/m²/年]',
        '  → 静的評価（GWP100）',
        '【DH2O】運用水消費量 [L/m²/日]',
        '',
        '■ 段階的seuil値（住宅集合住宅 Ic_construction）',
        '  Phase 1（2022-2024）: 640 kgCO₂eq/m²  ← 達成容易な基準から開始',
        '  Phase 2（2025-2027）: 530 kgCO₂eq/m²  (-17%)',
        '  Phase 3（2028予定）:  490 kgCO₂eq/m²  (-23%)',
        '  Phase 4（2031目標）:  350 kgCO₂eq/m²  (-45%)',
        '▶ 段階的引き下げ = 毎フェーズに技術革新プレッシャーを制度化',
    ])

    two_col(prs, 'Level(s) / EN15978 / EN15804+A2 との関係',
        'Level(s)・EN15978（整合）',
        ['・EU委員会策定の建物評価フレームワーク',
         '・Macro-obj 1.1 = WLC評価 = Ic指標に対応',
         '・EPBD 2024でLevel(s)活用を推奨',
         '・RE2020はLevel(s)の「先行実装」',
         '・計算ベース: EN 15978（建物LCA規格）',
         '',
         '▶ フランスがEU規格策定で主導権を持つ根拠'],
        'EN 15804+A2との矛盾（政治的争点）',
        ['・EN 15804+A2: 建材EPD欧州標準規格',
         '  バイオジェニック: −1/+1（時間評価なし）',
         '→ RE2020の動的LCAと手法的に非整合',
         '',
         '・CERIB等が「EU規格非整合」として批判',
         '・仏政府「EN 15804改訂で解消される」と反論',
         '',
         '・ISO 21391-1:2025（仏AFNOR主導）が',
         '  動的LCAの国際的調停役として機能'])

    # ── Section 2 ──────────────────────────────────────────────────────────────
    section(prs, '②', 'ダイナミックLCAの要件整理', 'No.9（ア）/ No.10（イ）/ No.11（ウ）')

    content(prs, 'DCF係数とFRE2020(t)の科学的根拠', [
        '■ なぜ将来排出のCO₂は「小さく」評価されるか',
        '  IPCC AR4 Bern2.5CCモデル: 放出CO₂の55%超が100年以内に海洋等に吸収',
        '  → 将来放出のCO₂は積分期間が短く累積放射強制力（AGWP）が小さい',
        '  DCF(t) = AGWP(t, 100年) / AGWP(0, 100年)',
        '',
        '■ 線形近似の導出（Benoist 2009 → Solinnen 2018）',
        '  IPCC AR4 Bern2.5CCのDCF(t)をt=0〜50年でR²=0.999で線形回帰',
        '  → FRE2020(t) = 1 − 0.00842 × t',
        '  t=50年: 線形=0.580 / Bern非線形=0.606（差4%）',
        '',
        '■ 制度化プロセス',
        '  2009: Benoist博士論文で原型  →  2018: Solinnen報告書で確定・Elodie実装',
        '  2021: RE2020 Arrêté（省令）で法規化',
        '',
        '■ 批判（Ventura, 2022）',
        '  固定DO=100年が t=50年排出を+25%過大評価',
        '  → Rivaton勧告No.5: 2028年Phase 3でISO 21391-1（変動DO）へ移行',
    ])

    content(prs, '動的LCA制度要件・段階施行スケジュール', [
        '■ FRE2020(t) 適用ルール（モジュール別）',
        '  A1-A3（t=0）:   F=1.000  ← 建設時排出・バイオ炭素固定',
        '  B段階（t≈25年）: F=0.790  ← 使用段階排出',
        '  C段階（t=50年）: F=0.580  ← 解体・廃棄段階排出',
        '  D（任意）:       F=0以下  ← Ic計算外・別途申告',
        '',
        '■ StockC（バイオジェニック炭素クレジット）',
        '  A1-A3: −Bc × 1.0  →  C4: +Bc × 0.58',
        '  正味 StockC = −0.42 × Bc [kgCO₂eq/m²]',
        '',
        '■ 主要パラメータ一覧',
        '  TOD（観測期間）:       100年（建物引渡しから、固定）',
        '  SLP（標準使用期間）:   50年（住宅）',
        '  係数下限:              0.58（t≥50年は0.58固定）',
        '  計算ツール:            Elodie（CSTB、無償、DHUP認定）',
        '  データベース:          INIES/FDES（約5,500件収録）',
        '',
        '■ 対象建物（段階施行）',
        '  2022〜: 住宅・オフィス・学校（新築）',
        '  2023〜: 商業施設・物流倉庫等に拡大',
    ])

    # ── Section 3 ──────────────────────────────────────────────────────────────
    section(prs, '③', '生物由来炭素・炭素貯蔵の制度的位置付け', 'No.12（ア）/ No.13（イ）/ No.14（ウ）')

    two_col(prs, 'StockC算定：静的LCAとの決定的な差',
        '静的LCA（EN 15804 −1/+1）',
        ['A1-A3（t=0）: −1.0 kgCO₂/kgC',
         'C4（t=50年）: +1.0 kgCO₂/kgC',
         '正味 = ゼロ（木材≡コンクリート）',
         '',
         '建材比較（静的LCA）:',
         '  CLT:      +98 kgCO₂eq/m³（排出）',
         '  RC造:    +300 kgCO₂eq/m³（排出）',
         '',
         '→ 設計者へのシグナル:',
         '  「木材もコンクリートも同じ（排出）」'],
        '動的LCA（RE2020 StockC）',
        ['A1-A3（t=0）: −1.0 × DCF(0)=1.0',
         'C4（t=50）:  +1.0 × DCF(50)=0.58',
         '正味 StockC = −0.42 kgCO₂/kgC',
         '',
         '建材比較（動的LCA）:',
         '  CLT:      −171 kgCO₂eq/m³（吸収！）',
         '  RC造:    +295 kgCO₂eq/m³（不変）',
         '  差: −269（CLT, 逆転）※数値要照合',
         '',
         '→ 設計者へのシグナル:',
         '  「木材は積極的炭素固定を意味する」'])

    content(prs, '評価期間（DO）の根拠・解体・再利用', [
        '■ なぜ「建設から100年」か（4理由）',
        '  ① GWP100との連続性（既存FDESデータ再利用が可能）',
        '  ② Levasseur (2010) との整合（動的LCA基礎論文がTOD=100年採用）',
        '  ③ 産業界への説明容易性（「同じ100年で精密に」）',
        '  ④ EU政策整合（EPBD・Level(s)もGWP100基準）',
        '',
        '■ 科学的批判と2028年改訂計画',
        '  Ventura (2022): 固定DO=100年 → t=50年排出を+25%過大評価',
        '  Rivaton No.5: Phase 3（2028）でISO 21391-1（変動DO）へ移行',
        '  Rivaton No.9: AR4→AR6パラメータ更新も同時実施',
        '',
        '■ 解体・再利用（C・Dモジュール）',
        '  C1-C4（t=50年）にF=0.58を適用。焼却→CO₂即時放出',
        '  D（任意）: Ic計算外。再利用で排出が後ろ倒し→追加クレジット可能性',
        '  → Rivaton No.10: Dモジュールの部分的Ic組込みを2028年〜試験実施',
        '',
        '■ SFM前提の条件（科学的留保）',
        '  StockCクレジットは「伐採＝再植林」（持続可能な森林管理）を前提',
        '  FSC/PEFC認証材でない場合の扱いは2028年改訂の議題',
    ])

    # ── Section 4 ──────────────────────────────────────────────────────────────
    section(prs, '④', '制度導入の背景および検討プロセス', 'No.4（ア）/ No.5（イ）/ No.6（ウ）/ No.7（エ）')

    content(prs, '制度化の4条件と15年の収束プロセス', [
        '■ 規制化の4条件（2021〜22年に同時成立）',
        '  ① 科学的基盤: Benoist 2009博論 → Solinnen 2018報告書で完成',
        '  ② データ基盤: E+C-プログラム（2016-19）で2,200棟実証・FDES整備',
        '  ③ 計算ツール: Elodie v2.5に動的LCA機能実装（2018）',
        '  ④ 政治的正当性: 産業界合意＋EUグリーンディール整合',
        '',
        '■ 15年の主要マイルストーン',
        '  2007: グルネル環境RTT — 炭素評価を政策課題化',
        '  2009: Benoist博論 — FRE2020(t)線形近似の原型',
        '  2016: E+C-開始（2,200棟実証）',
        '  2018: Solinnen報告書・Elodie実装完了',
        '  2019: EUグリーンディール — WLC義務化の政治的文脈確立',
        '  2021: RE2020法令公布（Décret 2021-1004）',
        '  2022: RE2020施行・EU議長国として国際発信',
        '  2025: ISO 21391-1:2025 策定（仏AFNOR主導）',
        '',
        '▶ E+C-プログラムの意義: 「実証なき規制化はしない」フランスの原則',
        '  2,200棟の実Ic分布データがPhase 1 seuil値（640）の根拠',
    ])

    two_col(prs, '所管省庁・組織体制とLoi Climat et Résilience',
        '主要機関と役割',
        ['【DHUP/DGALN】環境・住宅省',
         '  RE2020の立案・法制化を主導',
         '【CSTB】建築科学技術センター',
         '  Elodie開発・FICESデータベース運営',
         '【Cerema】環境・都市計画研究センター',
         '  技術ガイド作成・行政技術顧問',
         '【ADEME】環境・エネルギー管理庁',
         '  炭素評価手法の科学的レビュー',
         '【FCBA】仏木材・建設木材センター',
         '  動的LCA採用を科学的・政治的に支持',
         '  「正当な評価の回復」として産業訴求'],
        'Loi Climat et Résilience（Loi n°2021-1104）',
        ['2021年8月22日制定（気候・強靭性法）',
         '市民気候委員会（CCC）150提言を立法化',
         '',
         '建築関連条文:',
         '  第10条: 断熱性能段階的強化（Ic_énergie授権）',
         '  第14条: 建設GHG 2030年比50%削減目標',
         '  第49条: エンボディド炭素評価・開示義務',
         '  第181条: SNBC（国家低炭素戦略）の法律格上げ',
         '',
         'CCC提言S3.3（2020年6月）:',
         '「建設脱炭素化を2022年以前に開始」',
         '→ RE2020の2022年施行を加速',
         '▶ Loi Climat＝法的授権 / RE2020＝技術実装'])

    two_col(prs, '産業界の対立・調停と国際標準先行者戦略',
        '産業界の対立と政府の調停（3手法）',
        ['■ 反対派の論拠',
         '  CERIB（コンクリート研究所）:',
         '   「EN 15804/ISO 14067と整合しない」',
         '  FFA（仏鋼鉄連盟）:',
         '   「リサイクル鋼のDモジュールをIcに」',
         '',
         '■ 政府の調停メカニズム',
         '  ① 段階的施行: Phase1→4で産業界に猶予',
         '  ② 先送り: Dモジュール要求は2028年再検討',
         '  ③ 共同投資: FDES登録費用の国庫補助',
         '',
         '■ 代替案と却下理由（ch04 §4.1）',
         '  静的LCA継続: SNBC目標達成に不十分',
         '  固定補助金: 科学的根拠なし',
         '  完全非線形DCF: 2022年実装不可',
         '▶ 動的LCA線形近似 = 唯一の三要件充足手段'],
        '国際標準先行者戦略',
        ['① ISO 21391-1:2025 の主導',
         '  AFNORが策定をリード',
         '  RE2020の科学的手法が国際規格の基盤に',
         '',
         '② EPBD 2024への影響力',
         '  EU WLC義務化でRE2020が先行事例として参照',
         '',
         '③ EU議長国（2022年）の活用',
         '  RE2020施行と同時期に欧州議題設定を主導',
         '',
         '■ Rivaton (2025) 12項目勧告（主要）',
         '  No.1: FDESを10,000件へ（現5,500件）',
         '  No.5: Phase 3でISO 21391-1移行',
         '  No.6: デフォルト値見直し（2026中）',
         '  No.9: AR4→AR6パラメータ更新',
         '  No.10: Dモジュール部分的Ic試験（2028）'])

    # ── Section 5 ──────────────────────────────────────────────────────────────
    section(prs, '⑤', '日本への示唆・まとめ')

    content(prs, '日本への示唆：4条件の現状と優先課題', [
        '■ 日本の4条件現状（2026年）',
        '  ① 科学的基盤: ISO 21391-1:2025で国際基盤確立。JIS化が可能',
        '  ② データ基盤: IDEA時系列炭素フローデータが未整備（最重要課題）',
        '  ③ 計算ツール: 動的LCA対応国産ツールが存在しない（最優先開発）',
        '  ④ 政治的正当性: 木材利用促進基本計画・CN宣言と整合するが産業合意未形成',
        '',
        '■ 日本が避けるべきフランスの失敗',
        '  失敗①: DO固定（+25%過大評価） → 最初からISO 21391-1変動DOを採用',
        '  失敗②: データ後追い       → 規制化前にIDEA時系列データ整備を完了',
        '  失敗③: 中小業者コスト無視  → 計算ツール無償提供・行政補助を最初から設計',
        '',
        '■ 推奨ロードマップ（フランスのE+C-型）',
        '  2027〜2029: 実証プログラム設計・実施',
        '  2030〜2031: データ整備・ツール開発・人材育成',
        '  2032〜    : 段階的義務化（Phase 1相当）',
    ])

    content(prs, 'まとめ：調査ブロック①の主要知見', [
        '■ No.1-3・8（制度全体構造）',
        '  RE2020は3指標・4フェーズseuil・LCA義務化の世界初制度。',
        '  Level(s)・EN15978と整合。EN 15804+A2との矛盾は現在進行形の政治的争点。',
        '',
        '■ No.9-11（動的LCA要件）',
        '  FRE2020(t)=1−0.00842tはBenoist(2009)→Solinnen(2018)の10年検証の産物。',
        '  固定DO=100年の過大評価（+25%）問題は2028年Phase 3で対処予定。',
        '',
        '■ No.12-14（バイオジェニック炭素）',
        '  StockC=−0.42Bcが木材の「排出」→「吸収」逆転を生む（CLT: +98→−171）。',
        '  SFM前提・EN 15804矛盾・Dモジュール除外が残る構造的限界。',
        '',
        '■ No.4-7（制度導入背景）',
        '  4条件の2021〜22年収束が唯一の制度化要因。',
        '  Loi Climat（法的授権）＋産業界の段階的調停が政治的実現性を担保。',
        '  フランスのISO 21391-1主導は「自国手法の国際規格化」戦略。',
        '',
        '▶ 総括: 動的LCAの規制化は「木材優遇」でなく「測定ツールの修正」。',
        '   正確な計測が正確な設計シグナルを生み、正確なGHG削減につながる。',
    ])

    content(prs, '参考文献', [
        'Benoist, A. (2009). Thèse: carbone biogénique. Université Paris-Est.',
        'Solinnen (2018). Rapport DCF coefficients pour RE2020.',
        'Levasseur, A. et al. (2010). Considering time in LCA. Environ. Sci. Technol. 44(8):3169.',
        'Ventura, A. & Feraille, A. (2021). Dynamic LCA. Building and Environment 196:107777.',
        'Ventura, A. (2022). Conceptual issue of the dynamic GWP. Int. J. LCA.',
        'Shine, K.P. et al. (2005). Alternatives to GWP. Climatic Change 68(3):281.',
        'Forster, P. et al. (2007). IPCC AR4 WG1 Chapter 2. Cambridge Univ. Press.',
        'ISO 21391-1:2025. Dynamic LCA of biogenic carbon. ISO, Geneva.',
        'Cerema (2024). Guide technique RE2020. Cerema, Lyon.',
        'Rivaton, R. (2025). Bilan RE2020 et recommandations phase 3. DHUP/DGALN.',
        'DHUP (2021). Décret 2021-1004 et Arrêté du 4 août 2021. Journal Officiel.',
        'Loi n°2021-1104 du 22 août 2021. Journal officiel de la République française.',
        'CCC (2020). Propositions: Proposition S3.3. Convention Citoyenne pour le Climat.',
        'Hoxha, E. et al. (2020). Uncertainty in LCA. Building and Environment 175:106780.',
        'CERIB (2021). Analyse critique RE2020 GWP dynamique. (非公開資料)',
        'France Bois 2024 (2024). Plan France Bois 2024. Ministère de l\'Agriculture.',
    ])

    end(prs)

    os.makedirs(OUTDIR, exist_ok=True)
    prs.save(OUTPUT)
    print(f'✅ 保存完了: {OUTPUT}')
    print(f'   スライド数: {len(prs.slides)}')

if __name__ == '__main__':
    build()
