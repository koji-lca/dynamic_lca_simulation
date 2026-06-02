"""build_survey2_ppt.py — Survey 2 SuMPO PPTX builder (SKILL.md compliant)"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE = '/tmp/sumpo_skill/extracted/sumpo-ppt-v2/assets/template.pptx'
OUTDIR = os.path.join(os.path.dirname(__file__), 'slides')
OUTPUT  = os.path.join(OUTDIR, 'survey2_技術手法調査_SuMPO.pptx')

# SuMPO colors (SKILL.md)
CHATHAMS = RGBColor(0x12,0x54,0x86)
ABSOLUTE = RGBColor(0x00,0x57,0xBA)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
BLACK    = RGBColor(0x00,0x00,0x00)
DGRAY    = RGBColor(0x44,0x44,0x44)

L_COVER, L_SEC, L_CON, L_END = 0, 7, 8, 10
BL, BT, BW, BH = 510948, 1381566, 9721187, 5278273

def clear_slides(prs):
    R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    lst = prs.slides._sldIdLst
    for s in list(lst):
        rid = s.get(f'{{{R}}}id')
        if rid:
            try: prs.part.drop_rel(rid)
            except: pass
        lst.remove(s)

def set_title(slide, text, sz=20, color=CHATHAMS, bold=True):
    sh = slide.shapes.title
    if not sh: return
    sh.text_frame.clear()
    p = sh.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color

def body(slide, lines, l=BL, t=BT, w=BW, h=BH, sz=13):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        bold = line.startswith('■') or line.startswith('【') or line.startswith('▶')
        if line.startswith('■'):
            color = CHATHAMS
        elif line.startswith('▶') or line.startswith('【'):
            color = ABSOLUTE
        elif line.startswith('  '):
            color = DGRAY
        else:
            color = BLACK
        r = p.add_run(); r.text = line.lstrip()
        r.font.size = Pt(sz - (1 if line.startswith('  ') else 0))
        r.font.bold = bold; r.font.color.rgb = color
        p.space_after = Pt(1)

def cover(prs, title, sub, date):
    sl = prs.slides.add_slide(prs.slide_layouts[L_COVER])
    set_title(sl, title, sz=24, color=CHATHAMS)
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]; r = p.add_run()
            r.text = sub + '\n' + date; r.font.size = Pt(14); r.font.color.rgb = DGRAY

def section(prs, num, title, sub=''):
    sl = prs.slides.add_slide(prs.slide_layouts[L_SEC])
    set_title(sl, f'{num}　{title}', sz=26, color=WHITE)
    if sub:
        for ph in sl.placeholders:
            if ph.placeholder_format.idx not in (0, 8):
                ph.text = sub
                for p2 in ph.text_frame.paragraphs:
                    for r in p2.runs:
                        r.font.size = Pt(15); r.font.color.rgb = WHITE

def content(prs, title, lines, sz=13):
    sl = prs.slides.add_slide(prs.slide_layouts[L_CON])
    set_title(sl, title); body(sl, lines, sz=sz)

def two_col(prs, title, lt, ll, rt, rl):
    HW = (BW - 200000) // 2
    sl = prs.slides.add_slide(prs.slide_layouts[L_CON])
    set_title(sl, title)
    tb = sl.shapes.add_textbox(Emu(BL), Emu(BT), Emu(HW), Emu(280000))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = lt; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = CHATHAMS
    body(sl, ll, l=BL, t=BT+300000, w=HW, h=BH-300000, sz=12)
    RL = BL + HW + 200000
    tb2 = sl.shapes.add_textbox(Emu(RL), Emu(BT), Emu(HW), Emu(280000))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = rt; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = CHATHAMS
    body(sl, rl, l=RL, t=BT+300000, w=HW, h=BH-300000, sz=12)

def end(prs):
    prs.slides.add_slide(prs.slide_layouts[L_END])

# ── Slide content ──────────────────────────────────────────────────────────────
def build():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    cover(prs,
        '調査ブロック②　技術及び手法調査報告',
        '動的LCAの理論・算定方法・建築物LCAにおける位置付け',
        '一般社団法人サステナブル経営推進機構（SuMPO）　2026年6月')

    content(prs, '目次', [
        '⑤ 動的LCAの理論的枠組み（No.15〜17）',
        '  • 動的LCAの基本概念・静的LCAとの差異（No.15）',
        '  • 評価期間の科学的根拠と政治的均衡（No.16）',
        '  • 線形近似の精度・限界とISO 21391-1（No.17）',
        '',
        '⑥ 評価指標及び算定方法（No.18〜20）',
        '  • Ic_construction完全算定式・パラメータ根拠（No.18）',
        '  • 炭素循環（吸収・固定・放出）の扱い（No.19）',
        '  • 放射強制力の考え方との関係性（No.20）',
        '',
        '⑦ 建築物LCA算定における位置付け（No.21〜23）',
        '  • 算定環境（Elodie・INIES・人材育成）（No.21）',
        '  • 静的LCAとのハイブリッド設計関係性（No.22）',
        '  • 評価結果の利用方法・表示方法（No.23）',
        '',
        '⑧ 総合考察・日本への技術移転示唆',
        '  参考文献',
    ])

    # ── Section 5 ──────────────────────────────────────────────────────────────
    section(prs, '⑤', '動的LCAの理論的枠組み', 'No.15（ア）/ No.16（イ）/ No.17（ウ）')

    content(prs, '動的LCAの基本概念：静的LCAが見落とすもの（No.15）', [
        '■ 静的LCA（GWP100）の構造的限界',
        '  「建物ライフサイクルのすべての排出・吸収を同一時点で評価する」という前提',
        '  → 木材：t=0に吸収、t=50年に放出 → 静的LCAでは正味ゼロ（実態と乖離）',
        '  → コンクリート：製造時排出のみ → 木材と同じ「炭素中立」の誤解を生む',
        '',
        '■ 動的LCAの2つの拡張',
        '【動的LCI】排出・吸収の発生タイミングを時系列で記録',
        '  各フロー（材料採取・製造・使用・廃棄）に「いつ起きるか」を付与',
        '【動的LCIA】時間依存型特性化係数（DCF）を適用',
        '  DCF(t) = 放射強制力の時間積分比率（後述）',
        '',
        '■ DCF(t)の直感的理解',
        '  今日放出した1kgCO₂は100年間大気に影響 → DCF=1.0（基準）',
        '  50年後に放出した1kgCO₂は残り50年しか影響しない → DCF=0.58',
        '  → 「後に放出されるほど気候への悪影響が小さい」',
        '',
        '▶ RE2020は動的LCAを「静的LCA＋後付け時間係数」として簡略実装',
        '   FRE2020(t) = 1 − 0.00842×t（線形近似）',
    ])

    two_col(prs, '評価期間の科学的根拠と政治的均衡（No.16）',
        '「100年」の科学的根拠',
        ['■ IPCC AR4 Bern2.5CC応答関数',
         '  C(t) = 0.217 + 0.259·e^(−t/172.9)',
         '       + 0.338·e^(−t/18.51)',
         '       + 0.186·e^(−t/1.186)',
         '',
         '  約22%（a₀=0.217）は事実上永続',
         '  → 100年間を超えても気候影響が続く',
         '',
         'IPCC AR5（1990）の政策的経緯:',
         '  GWP20/GWP100/GWP500 を提示',
         '  政策決定でGWP100が選択された背景:',
         '  「短期: メタン有利、長期: CO₂有利」',
         '  → 気候科学者・政策担当者が合意した',
         '    政治的・実用的な均衡点'],
        '「100年」への科学的批判',
        ['■ Ventura (2022) の批判',
         '  固定DO（観測窓）=100年 vs',
         '  変動DO（ISO 21391-1方式）の差',
         '',
         '  固定DO=100年では t=50年排出を',
         '  +25%過大評価（Ventura, 2022）',
         '',
         '■ Shine et al. (2005) の代替案',
         '  GTP（地球温度変動ポテンシャル）',
         '  提案。IGWPとの差を定量化。',
         '  → RE2020採用なし（データ要求過大）',
         '',
         '■ 2028年改訂（Rivaton No.5）',
         '  ISO 21391-1:2025の変動DO方式へ',
         '  移行検討。AR6パラメータ更新も同時',
         '',
         '▶ 100年は科学的最適ではなく',
         '   政治的均衡点として機能'])

    content(prs, '線形近似の精度・限界とISO 21391-1（No.17）', [
        '■ 線形近似 FRE2020(t) の精度',
        '  BernモデルDCF(t)をt=0〜50でR²=0.9997で線形近似（Benoist, 2009）',
        '  t=50年時の誤差: |0.580 − 0.606| / 0.606 ≈ 4.3%（実用上十分）',
        '',
        '■ 線形近似の3つの限界',
        '  ① t=50年以降は DCF実値がFRE2020を上回る（線形が過小評価）',
        '  ② CH₄・N₂OにはCO₂とは異なるDCF(t)曲線が必要（RE2020は未対応）',
        '  ③ Ventura(2022)の固定DO=100年問題（+25%過大評価）',
        '',
        '■ ISO 21391-1:2025（仏AFNOR主導、2025年策定）との比較',
        '  RE2020:                 ISO 21391-1:2025',
        '  FRE2020(t)=1−0.00842t  非線形DCF（Bern数値積分直接使用）',
        '  固定DO=100年            変動DO（排出時点から100年）',
        '  CO₂のみ動的評価         各ガス種別に動的評価',
        '  線形近似（実装容易）     高精度（計算負荷増）',
        '',
        '▶ RE2020の線形近似は「過渡期の簡易化」としてISO 21391-1に位置付けられる',
        '   2028年Phase 3で非線形移行を検討中',
    ])

    # ── Section 6 ──────────────────────────────────────────────────────────────
    section(prs, '⑥', '評価指標及び算定方法', 'No.18（ア）/ No.19（イ）/ No.20（ウ）')

    content(prs, 'Ic_construction 完全算定式とパラメータ根拠（No.18）', [
        '■ Ic_construction の完全算定式',
        '  Ic = (1/SHON) × Σ_m [ Σ_i E_i^(m) × GWP₁₀₀(i) × FRE2020(t_m) ] + StockC',
        '',
        '  SHON: 正味延べ面積 [m²]',
        '  E_i^(m): モジュールm における物質iの排出量 [kg]',
        '  GWP₁₀₀(i): 静的温暖化係数（IPCC AR4）',
        '  FRE2020(t_m): モジュールm の時間係数',
        '  StockC: バイオジェニック炭素クレジット（負値）',
        '',
        '■ モジュール別時間係数',
        '  A1-A3（t=0年）:  FRE2020 = 1.000  ← 建設時',
        '  B（t≈25年）:    FRE2020 = 0.790  ← 使用段階平均',
        '  C（t=50年）:    FRE2020 = 0.580  ← 解体・廃棄',
        '',
        '■ 係数0.00842の導出プロセス（2009→2021）',
        '  Step 1: AR4 Bern2.5CC応答関数でCO₂の大気中残存率を計算',
        '  Step 2: AGWP(t, 100)を数値積分',
        '  Step 3: DCF(t) = AGWP(t,100) / AGWP(0,100)',
        '  Step 4: t=0〜50年で線形回帰 → 傾き = −0.00842（R²=0.9997）',
        '  → 2021年 Arrêté（省令）で法規化',
    ])

    two_col(prs, '炭素循環（吸収・固定・放出）の扱い（No.19）',
        '木材のバイオジェニック炭素サイクル',
        ['大気CO₂',
         '↓ 光合成（NPP: 0.5〜3.0 kgC/m²/年）',
         '樹木バイオマス（生体）',
         '↓ 伐採・搬出',
         '製材・CLT（製品内貯蔵）',
         '↓ 建物として使用（50年以上）',
         '建物構造材（長期貯蔵）',
         '↓ 解体後',
         '  ① 焼却（C4）→ CO₂即時放出',
         '  ② 埋立（C4）→ CH₄+CO₂ゆっくり放出',
         '  ③ 再利用（D）→ さらに貯蔵延長',
         '大気CO₂（帰還）',
         '',
         '▶ RE2020: ③再利用はIc計算外（Dモジュール）'],
        'SFM前提・EN 16485との関係',
        ['■ SFM（持続可能な森林管理）前提',
         '  クレジットは「伐採=再植林」が条件',
         '  FSC/PEFC認証材が対象',
         '',
         '  リスク（RE2020の潜在的課題）:',
         '  ・非認証材（海外輸入材等）',
         '  ・気候変動による森林劣化',
         '  ・伐採直後の一時的炭素負債',
         '  → Rivaton No.7: 対処を2028年検討',
         '',
         '■ IPCC HWP会計との関係',
         '  建築用途構造材（半減期50年以上）',
         '  → RE2020 StockC計算の対象',
         '',
         '■ EN 16485:2014との差異',
         '  EN 16485: 時間軸評価なし（静的）',
         '  RE2020: FRE2020(t)適用・StockCクレジット認める',
         '→ 欧州規格との方法論的差異が継続'])

    content(prs, '放射強制力の考え方との関係性（No.20）', [
        '■ 放射強制力（RF）とは',
        '  大気成分の変化による地球エネルギー収支の擾乱 [W/m²]',
        '  RF_CO₂(C) = 5.35 × ln(C/C₀)  [W/m²]',
        '  C₀: 産業化前CO₂濃度（278 ppm）',
        '',
        '■ GWP100とAGWPの関係',
        '  GWP_H = AGWP_H(x) / AGWP_H(CO₂)',
        '  AGWP: 累積放射強制力（時間積分）',
        '  AGWP₁₀₀(CO₂) ≈ 8.69 × 10⁻¹⁴ W m⁻² yr kg⁻¹（AR4）',
        '',
        '■ なぜ「将来放出のCO₂は気候影響が小さい」か',
        '  観測窓（TOD=100年）の終点に近い時点で放出されたCO₂は',
        '  積分期間が残り少ない → AGWP(t,100)が小さくなる',
        '  DCF(50) = AGWP(50,100) / AGWP(0,100) ≈ 0.58',
        '  「今日の1kgCO₂」vs「50年後の1kgCO₂」の気候影響比 = 100:58',
        '',
        '■ IPCC AR6（2021）への更新影響',
        '  AR6: a₀=0.200（AR4: 0.217）、τ₁=394.4年（AR4: 172.9年）',
        '  → DCF(t)がやや大きくなる（木材に若干不利）',
        '  → RE2020は現在もAR4使用。AR6更新は2028年Phase 3の議題（Rivaton No.9）',
        '',
        '▶ 方法論の正当性はAR4→AR6更新で変わらない。パラメータ値が微調整されるのみ。',
    ])

    # ── Section 7 ──────────────────────────────────────────────────────────────
    section(prs, '⑦', '建築物LCA算定における位置付け', 'No.21（ア）/ No.22（イ）/ No.23（ウ）')

    content(prs, 'フランスのLCA算定エコシステム（No.21）', [
        '■ 三層構造',
        '【第1層】データベース: INIES（inies.fr）',
        '  FDES（建材環境製品宣言）約5,500件収録（2026年現在）',
        '  デフォルト値: FDES未整備建材用「保守的（不利側）」標準値',
        '',
        '【第2層】計算ソフト: Elodie（CSTB開発・無償提供）',
        '  FDESをINIESから直接インポート・Ic計算自動化',
        '  動的LCA: FRE2020(t)の自動適用・StockC自動算定',
        '  出力: 申請用PDFレポート（DHUP指定フォーマット）',
        '  精度: 整合化FDES使用時 ±5%以内（Cerema, 2024 §6.2）',
        '',
        '【第3層】申請・認証機関',
        '  建築確認申請（Permis de construire）に計算書を添付（法的義務）',
        '  任意認証: BBCA（seuil×0.70以下）、NF HQE等',
        '',
        '■ 日本との比較（計算ツール・データベース件数）',
        '  フランス（Elodie、INIES）:   5,500件  / 動的LCA対応 ✅',
        '  ドイツ（ÖKOBAUDAT準拠）:    1,200件  / 動的LCA対応 ✗',
        '  日本（IDEA-BIM）:           3,000件  / 動的LCA対応 ✗  ← 最優先課題',
    ])

    two_col(prs, '静的LCAとのハイブリッド設計（No.22）',
        'RE2020のハイブリッド設計思想',
        ['純粋な動的LCAではなく:',
         '「静的LCA結果 × 動的係数（後付け）」',
         '',
         '静的LCAデータ（FDES/EPD）',
         '  × FRE2020(t_m)',
         '+ StockC（バイオジェニック）',
         '= Ic_construction',
         '',
         '■ この方式の合理性',
         '・既存FDESインフラを流用（コスト最小化）',
         '・動的LCI整備が不要',
         '・Elodieへの実装が「係数の掛け算のみ」',
         '',
         '■ 差異が大きいケース（動的補正が効く）',
         '  CLT・LVL等バイオジェニック炭素多',
         '  B段階で大量修繕が発生する建材',
         '',
         '■ 差異が小さいケース',
         '  コンクリート・鋼材（バイオ炭素なし）'],
        'EN 15804+A2との方法論的矛盾',
        ['| 項目 | EN 15804+A2 | RE2020 |',
         '|------|-------------|--------|',
         '| 時間軸評価 | なし | あり |',
         '| バイオ炭素 | −1/+1正味ゼロ | StockCクレジット |',
         '| 評価基準 | 欧州委員会策定 | 仏国内法規 |',
         '',
         '■ 矛盾の政治的重要性',
         '  CERIB: 「EU規格非整合」として継続批判',
         '  フランス政府: 「EN 15804改訂で解消」と反論',
         '',
         '▶ ISO 21391-1:2025が国際的調停役',
         '   → 将来的にEN 15804改訂の動力となる可能性'])

    content(prs, '評価結果の利用方法・表示方法（No.23）', [
        '■ Ic指標の法的開示フロー',
        '  ① 設計者がElodieでIc試算 → 設計段階での材料比較・最適化',
        '  ② 建築確認申請にElodie出力PDFを添付（法的義務）',
        '  ③ 施工段階: 設計LCA→実施LCAへ更新（A5実績反映）',
        '  ④ 完成時: 最終Ic値を確定・認証機関へ提出（記録10年保存）',
        '',
        '■ 申告書フォーマット（例）',
        '  Ic_construction: 412 kgCO₂eq/m²  ≤ seuil 640  ✅',
        '    うち StockC:    −85 kgCO₂eq/m²（バイオジェニック炭素クレジット）',
        '  Ic_énergie:      2.9 kgCO₂eq/m²/年 ≤ seuil 4.0  ✅',
        '',
        '■ 金融・保険分野での活用',
        '  グリーンボンド: EUタクソノミー技術基準の参照指標',
        '  住宅ローン: BBCA認証建物への金利優遇（Crédit Agricole等）',
        '  ESG報告: 不動産ポートフォリオのScope 3計算にIc値を活用',
        '',
        '■ 国際比較（建物LCA開示義務）',
        '  フランス（RE2020、2022〜）: 義務・動的LCAあり',
        '  スウェーデン（PBL、2022〜）: 義務・動的LCAなし',
        '  EU全体（EPBD 2024）:        2030年目標・義務化予定',
        '  日本（CASBEE等）:            義務なし・最優先整備課題',
    ])

    # ── Section 8 ──────────────────────────────────────────────────────────────
    section(prs, '⑧', '総合考察・日本への技術移転示唆')

    content(prs, '動的LCA手法比較：RE2020 vs ISO 21391-1 vs 純粋動的LCA', [
        '■ 3手法の方法論的比較',
        '  手法              | DCF計算    | DO設定    | 実装容易性 | データ要件',
        '  RE2020（現行）     | 線形近似   | 固定100年 | 高（簡単） | 低（既存FDES）',
        '  ISO 21391-1:2025  | 非線形Bern | 変動DO    | 中         | 中',
        '  純粋動的LCA       | 数値積分   | 変動DO    | 低（複雑） | 高（時系列LCI）',
        '',
        '■ 各手法の政策適用可能性',
        '  RE2020方式: 即時実施可。既存LCAデータ（静的EPD）を活用。「過渡期の手法」',
        '  ISO 21391-1: 精度と実用性の次世代バランス。日本のJIS化候補',
        '  純粋動的LCA: 研究・高精度評価向け。義務規制には不向き（現時点）',
        '',
        '■ 日本が「スキップ」できる失敗',
        '  失敗①: 固定DO採用 → ISO 21391-1（変動DO）を最初から採用すればよい',
        '  失敗②: データ後追い → 規制化前にIDEA時系列LCIデータを整備',
        '  失敗③: 単一係数（0.00842）の固定化 → AR6ベースの係数から出発',
        '',
        '▶ 日本はフランスの「過渡期」を飛ばし、ISO 21391-1方式で直接スタートできる',
        '   → ただしデータ整備（IDEA時系列化）と計算ツール開発が必要条件',
    ])

    content(prs, 'まとめ：調査ブロック②の主要知見（No.15〜23全カバー）', [
        '■ No.15〜17（動的LCA理論的枠組み）',
        '  動的LCAは「静的LCA＋時間軸」の拡張。DCF(t)は物理的放射強制力に根拠。',
        '  評価期間100年は科学的最適でなく政治的均衡点。線形近似は精度95%超（4%誤差）。',
        '  ISO 21391-1:2025が次世代標準として非線形・変動DOを国際的に定着させた。',
        '',
        '■ No.18〜20（評価指標・算定方法）',
        '  係数0.00842はBenoist(2009)→Solinnen(2018)の10年検証→2021年法規化の産物。',
        '  炭素循環は「光合成→製品内貯蔵→廃棄放出」の全フローをDCF適用で追跡。',
        '  放射強制力ベースの動的評価は「同じ1kgでも排出時点が違えば影響が違う」を正確に計測。',
        '',
        '■ No.21〜23（建築物LCA算定の位置付け）',
        '  フランスの三層エコシステム（INIES→Elodie→申請）は無償・法的認定が特徴。',
        '  RE2020は「静的LCA×動的係数」のハイブリッド。既存インフラ活用が鍵。',
        '  Ic指標はグリーンボンド・住宅ローン・ESG報告へと金融分野にも展開中。',
        '',
        '▶ 技術移転の優先順位: ①IDEA時系列LCI整備  ②動的LCA計算ツール開発',
        '                        ③ISO 21391-1のJIS化 ④実証プログラム（E+C-型）',
    ])

    content(prs, '参考文献', [
        'Levasseur, A. et al. (2010). Considering time in LCA. Environ. Sci. Technol. 44(8):3169.',
        'Levasseur, A. et al. (2013). Biogenic carbon & dynamic LCA. J. Industrial Ecology 17(1):117.',
        'Ventura, A. & Feraille, A. (2021). Dynamic LCA. Building and Environment 196:107777.',
        'Ventura, A. (2022). Conceptual issue of the dynamic GWP. Int. J. LCA.',
        'Benoist, A. (2009). Thèse: carbone biogénique. Université Paris-Est.',
        'Solinnen (2018). Rapport DCF coefficients pour RE2020.',
        'Forster, P. et al. (2007). IPCC AR4 WG1 Chapter 2. Cambridge Univ. Press.',
        'IPCC (2021). AR6 WG1 Chapter 7: The Earth\'s Energy Budget.',
        'Eggleston, H.S. et al. (2006). IPCC 2006 GL: Vol.4 AFOLU. IPCC, Geneva.',
        'Shine, K.P. et al. (2005). Alternatives to GWP. Climatic Change 68(3):281.',
        'EN 16485:2014. Round and sawn timber — Environmental product declarations. CEN.',
        'ISO 21391-1:2025. Dynamic LCA of biogenic carbon. ISO, Geneva.',
        'CSTB (2022). Elodie: Guide d\'utilisation v4.5. CSTB, Paris.',
        'Cerema (2024). Guide technique RE2020. Cerema, Lyon.',
        'Rivaton, R. (2025). Bilan RE2020 et recommandations phase 3. DHUP/DGALN.',
        'INIES (2024). Base de données INIES: mode d\'emploi 2024. CSTB/FFB.',
        'Houghton, R.A. (2007). Balancing the global carbon budget. Ann. Rev. Earth Planet. Sci. 35:313.',
        'Grubb, M. et al. (2021). Planetary Economics. Earthscan/Routledge.',
    ])

    end(prs)

    os.makedirs(OUTDIR, exist_ok=True)
    prs.save(OUTPUT)
    print(f'✅ 保存完了: {OUTPUT}')
    print(f'   スライド数: {len(prs.slides)}')

if __name__ == '__main__':
    build()
