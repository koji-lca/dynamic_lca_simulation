#!/usr/bin/env python3
"""
SuMPO 動的LCA調査 統合プレゼンテーション生成スクリプト
調査ブロック①（政策・制度）+ 調査ブロック②（技術・手法）を一本化。
HTMLレポートの内容をベースに重複除去・冗長表現を排除して構成。
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ── Paths ──────────────────────────────────────────────────────────────
HERE = Path(__file__).parent
TEMPLATE = HERE.parent / "shared" / "templates" / "sumpo_template.pptx"
CHARTS   = HERE.parent / "deliverables" / "html_report" / "assets" / "charts"
OUTPUT   = HERE / "SuMPO_動的LCA調査_統合報告書.pptx"

# ── Brand colors ───────────────────────────────────────────────────────
C_TITLE  = RGBColor(0x12, 0x54, 0x86)
C_DARK   = RGBColor(0x49, 0x4D, 0x51)
C_SILVER = RGBColor(0xC5, 0xC9, 0xCD)
C_NAVY   = RGBColor(0x00, 0x57, 0xBA)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xE0, 0x6C, 0x00)
C_GREEN  = RGBColor(0x1A, 0x8A, 0x3A)

# ── Layout indices ─────────────────────────────────────────────────────
L_COVER   = 0
L_SEC_BLU = 1
L_CON_BLU = 2
L_SEC_NAV = 7
L_CON_NAV = 8
L_END_BLU = 10


# ── Helpers ────────────────────────────────────────────────────────────

def _add_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _set_title(slide, text, size=22):
    sh = slide.shapes.title
    if sh is None:
        return
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = C_TITLE


def _fill_body(slide, items, size=15, ph_idx=1):
    """
    items: str → level 0
           (str, level) → specified level
           '' → blank line
    """
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in items:
        text  = item[0] if isinstance(item, tuple) else item
        level = item[1] if isinstance(item, tuple) else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text  = text
        p.level = level
        if p.runs:
            p.runs[0].font.size = Pt(size)
        else:
            p.font.size = Pt(size)


def _add_table(slide, headers, rows, x, y, w, h, hsize=12, bsize=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        tf = cell.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = hdr
        run.font.size = Pt(hsize)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_SILVER
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = val
            run.font.size = Pt(bsize)


def _add_chart(slide, fname, x, y, w, h):
    img = CHARTS / fname
    if img.exists():
        slide.shapes.add_picture(str(img), x, y, w, h)


def _key_box(slide, text, x=None, y=None, w=None, h=None, size=13):
    x = x or Emu(400_000)
    y = y or Emu(5_600_000)
    w = w or Emu(9_900_000)
    h = h or Emu(900_000)
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = C_ORANGE


def _set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── Slide builders ─────────────────────────────────────────────────────

def s01_cover(prs):
    slide = _add_slide(prs, L_COVER)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "フランスRE2020における動的LCA\n政策・技術統合調査報告書"
        run.font.size = Pt(26)
        run.font.bold = True
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            lines = [
                "調査ブロック① 政策及び制度調査　＋　調査ブロック② 技術及び手法調査",
                "",
                "一般社団法人サステナブル経営推進機構（SuMPO）",
                "2026年6月",
            ]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
    return slide


def s02_agenda(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "目次")
    _fill_body(slide, [
        "Ⅰ．なぜ動的LCAか ── 静的GWP100の構造的欠陥",
        "Ⅱ．フランスの制度化プロセス ── 15年の収束",
        "Ⅲ．RE2020の制度構造 ── 3指標・seuil値・段階施行",
        "Ⅳ．動的LCAの理論 ── DCF・Bern2.5CCモデル・0.00842係数",
        "Ⅴ．バイオジェニック炭素とStockC ── 算定方法と数値例",
        "Ⅵ．評価期間DO=100年 ── 根拠・批判・ISO 21391-1",
        "Ⅶ．算定エコシステム ── Elodie・ハイブリッド設計・開示制度",
        "Ⅷ．現行制度の限界と日本への示唆",
    ], size=16)
    return slide


# ── Section Ⅰ ──────────────────────────────────────────────────────────

def s03_sec_why(prs):
    slide = _add_slide(prs, L_SEC_NAV)
    _set_title(slide, "Ⅰ．なぜ動的LCAか", size=28)
    _fill_body(slide, ["静的GWP100が「物理的に誤った答え」を出し続けた"], size=18)
    return slide


def s04_gwp100_failure(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "静的GWP100の構造的欠陥")
    _fill_body(slide, [
        "【前提の問題】静的LCAは「今日の1 kgCO₂ ≡ 50年後の1 kgCO₂」と仮定する",
        "",
        "木材建材（CLT）のライフサイクル:",
        ("A1-A3（建設 t=0）: バイオジェニック炭素固定 → −269 kgCO₂eq/m³", 1),
        ("C4（解体 t=50年）: 炭素放出      → +269 kgCO₂eq/m³", 1),
        ("静的LCA合計 = 0（木材 ≡ コンクリート）← 物理的に誤り", 1),
        "",
        "【物理的現実】50年後の1 kgCO₂排出は、今日の排出より気候影響が42%小さい",
        ("根拠: IPCC AR4 Bern2.5CCモデル — 海洋・生物圏が50年間CO₂を吸収し続ける", 1),
        "",
        "動的LCA（RE2020）での評価:",
        ("A1-A3: −269（即時クレジット）　C4: +269×0.58 = +156", 1),
        ("正味 StockC = −113 kgCO₂eq/m³ ← 実在する気候便益を初めて計測", 1),
    ], size=14)
    _key_box(slide, "動的LCA導入は「木材優遇」ではなく「測定ツールの誤り修正」である")
    _add_chart(slide, "fig07_module_dcf_impact.png",
               Emu(6_500_000), Emu(1_600_000), Emu(3_200_000), Emu(3_200_000))
    _set_notes(slide, "CLT静的+98 → 動的−171 kgCO₂eq/m³。−269の逆転が生じる。")
    return slide


def s05_4conditions(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "フランスが2022年に規制化できた4条件の収束")
    _add_table(slide,
        ["条件", "必要な状態", "フランスの達成時期"],
        [
            ["① 科学的基盤", "査読を経た係数確立", "2009年 Benoist → 2018年 Solinnen"],
            ["② データ基盤", "十分なFDES/EPDカバレッジ", "E+C-実証（2016-2019）で整備"],
            ["③ 計算ツール", "実務者が使えるソフトウェア", "Elodie が動的LCA計算を自動化（2018〜）"],
            ["④ 政治的正当性", "産業界・規制当局の合意", "E+C-データ＋FCBA支持＋EUグリーンディール"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(2_400_000),
        hsize=13, bsize=12
    )
    _add_chart(slide, "fig12_policy_timeline.png",
               Emu(400_000), Emu(4_200_000), Emu(9_900_000), Emu(2_400_000))
    _set_notes(slide, "Benoist 2009博士論文 → Solinnen 2018実務化 → E+C- 2,200棟実証 → 2022施行。10年の準備を省略しなかった。")
    return slide


# ── Section Ⅱ ──────────────────────────────────────────────────────────

def s06_sec_re2020(prs):
    slide = _add_slide(prs, L_SEC_BLU)
    _set_title(slide, "Ⅱ．フランスRE2020の制度構造", size=28)
    _fill_body(slide, ["3指標体系・seuil値・4段階施行の概要"], size=18)
    return slide


def s07_3indicators(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "RE2020の3評価指標と炭素指標 Ic の構造")
    _add_table(slide,
        ["指標", "単位", "評価対象", "RT2012との関係"],
        [
            ["Bbio", "無次元", "断熱・日照・通風設計（冷暖房需要）", "基準値を約30%強化"],
            ["Cep,nr", "kWh/m²/yr", "非再生可能一次エネルギー消費（運用段階）", "継続指標"],
            ["Ic", "kgCO₂eq/m²", "ライフサイクルGHG（動的LCA補正込み）", "新設 ― 世界初の法定義務化"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(1_900_000),
        hsize=13, bsize=12
    )
    _fill_body(slide, [
        "Ic の構成式:   Ic = Ic_construction + Ic_énergie − StockC",
        "",
        ("Ic_construction: 建材・施工段階（A1-A5, B1-B5, C1-C4）の炭素排出量", 1),
        ("Ic_énergie:     運用エネルギー（B6）の炭素排出量", 1),
        ("StockC:         バイオジェニック炭素貯蔵クレジット（動的補正込み）← RE2020独自拡張", 1),
        "",
        "評価根拠: EN 15978モジュール構造（A1〜D）を採用。StockCはRE2020が追加したモジュール。",
    ], size=13, ph_idx=1)
    _set_notes(slide, "3指標の相補性: 省エネ→断熱材増量→エンボディド炭素増加というトレードオフをIcが可視化。")
    return slide


def s08_seuil(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "段階的施行スケジュールと Ic_construction 上限値（seuil）")
    _add_table(slide,
        ["建物種別", "Phase 1（2022）", "Phase 2（2025）\n−15%", "Phase 3（2028）\n−25%", "Phase 4（2031）\n−30〜40%"],
        [
            ["戸建住宅", "640", "約545", "約480", "約420〜448"],
            ["集合住宅", "840", "約715", "約630", "約505〜588"],
            ["事務所",   "740", "約630", "約555", "約444〜518"],
            ["教育施設", "660", "約560", "約495", "約396〜462"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(5_500_000), Emu(2_300_000),
        hsize=12, bsize=11
    )
    _add_chart(slide, "fig08_seuil_timeline.png",
               Emu(6_100_000), Emu(1_700_000), Emu(4_200_000), Emu(3_000_000))
    _fill_body(slide, [
        "単位: kgCO₂eq/m²（代表値: H2b気候ゾーン・標準規模。実際は気候ゾーン・面積・階数で細分化）",
        "seuil値はE+C-実証2,200棟データから導出。根拠ある数値であることが産業界の合意を可能にした。",
    ], size=12, ph_idx=1)
    _set_notes(slide, "集合住宅（840）が戸建（640）より高い理由: 共用設備・地下駐車場の炭素排出が大きいため。")
    return slide


# ── Section Ⅲ ──────────────────────────────────────────────────────────

def s09_sec_theory(prs):
    slide = _add_slide(prs, L_SEC_NAV)
    _set_title(slide, "Ⅲ．動的LCAの理論的枠組み", size=28)
    _fill_body(slide, ["DCF・Bern2.5CCモデル・0.00842係数の科学的根拠"], size=18)
    return slide


def s10_dcf_bern(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "動的特性化係数（DCF）とIPCC AR4 Bern2.5CCモデル")
    _fill_body(slide, [
        "DCF の定義（Levasseur et al., 2010）:",
        ("DCF(t_i, DO) = AGWP(t_i, DO) / AGWP(0, DO)", 1),
        ("AGWP = 時刻 t_i に排出した1 kgCO₂が DO 年間に生じる累積放射強制力", 1),
        ("物理的意味: 早期排出ほど1に近く（影響大）、後年排出ほど0に近い（影響小）", 1),
        "",
        "CO₂大気残存関数（IPCC AR4 Bern2.5CC）:  r(t) = 0.217 + Σ c_j·exp(−t/τ_j)",
        ("c₁=0.259, τ₁=172.9年（深層海洋） | c₂=0.338, τ₂=18.51年（中間層）", 1),
        ("c₃=0.186, τ₃=1.186年（生物圏・土壌） | a₀=0.217（永続残存）", 1),
        "",
        "DO=100年での典型値:",
        ("t=0年: DCF=1.000 | t=25年: DCF=0.792 | t=50年: DCF=0.606（≈42%減）", 1),
    ], size=13)
    _add_chart(slide, "fig01_bern_impulse.png",
               Emu(6_600_000), Emu(1_800_000), Emu(3_100_000), Emu(2_600_000))
    _add_chart(slide, "fig02_crf_dcf_relationship.png",
               Emu(6_600_000), Emu(4_500_000), Emu(3_100_000), Emu(2_100_000))
    _set_notes(slide, "r(t)の4指数モデルは「深海/中間層/生物圏+永続分」の物理吸収メカニズムを表現。放射強制力係数 a_x = 1.55×10⁻¹⁵ W·m⁻²·(kgCO₂)⁻¹。")
    return slide


def s11_0842(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "RE2020線形近似係数 0.00842 の導出と精度")
    _fill_body(slide, [
        "導出プロセス（Benoist, 2009 → Solinnen, 2018）:",
        ("Step 1: Bern2.5CCモデルでAGWP(t)を数値積分 → DCF(t)を計算", 1),
        ("Step 2: t=0〜50年の範囲を最小二乗法で線形回帰", 1),
        ("結果: F_RE2020(t) = 1 − 0.00842·t（R²=0.9997）", 1),
        ("下限: t>50年では0.580固定（建物LCのEoLは通常t=50年）", 1),
        "",
        "精度（Bern非線形DCFとの比較）:",
        ("t=10年: 誤差+0.1%　t=25年: −0.3%　t=50年: −4.3%（過小評価）", 1),
        "",
        "なぜ線形近似か:",
        ("透明性・実務コスト・規制予測可能性を優先。2022年段階では最適な選択", 1),
        ("2028年フェーズ3でISO 21391-1:2025（非線形DCF）への移行を検討中", 1),
    ], size=14)
    _add_chart(slide, "fig03_dcf_re2020_vs_bern.png",
               Emu(6_100_000), Emu(1_700_000), Emu(4_200_000), Emu(3_400_000))
    _key_box(slide, "係数は10年の段階的検証（学術→実証→実務化→法規）を経て確定")
    return slide


# ── Section Ⅳ ──────────────────────────────────────────────────────────

def s12_sec_biogenic(prs):
    slide = _add_slide(prs, L_SEC_BLU)
    _set_title(slide, "Ⅳ．バイオジェニック炭素とStockC", size=28)
    _fill_body(slide, ["算定方法・数値例・持続可能な森林管理（SFM）前提"], size=18)
    return slide


def s13_static_vs_dynamic(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "静的 vs 動的LCA：建材評価の逆転")
    _add_table(slide,
        ["建材", "静的LCA [kgCO₂eq/m³]", "動的LCA [kgCO₂eq/m³]", "変化"],
        [
            ["CLT（直交集成板）", "+98（排出）", "−171（吸収）", "−269 逆転"],
            ["製材",             "+65（排出）", "−115（吸収）", "−180 逆転"],
            ["RC造コンクリート", "+300前後",    "+295前後",     "ほぼ不変"],
            ["構造用鋼材",       "+550前後",    "+548前後",     "ほぼ不変"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(6_000_000), Emu(2_200_000),
        hsize=13, bsize=12
    )
    _fill_body(slide, [
        "逆転のメカニズム: 木材は「今固定（t=0, DCF=1.0）・50年後放出（DCF=0.58）」",
        "コンクリート・鉄鋼は化石由来排出のみで、時間補正の影響はほぼゼロ",
        "",
        "政策的含意: 静的LCAを使い続ける限り、設計者は「木材もコンクリートも同じ」という",
        "誤ったシグナルを受け取る。動的LCAへの移行は補助金でなく正確な計測への移行。",
    ], size=13, ph_idx=1)
    _add_chart(slide, "fig04_wood_gwp_comparison.png",
               Emu(6_200_000), Emu(3_900_000), Emu(3_600_000), Emu(2_700_000))
    return slide


def s14_stockc_calc(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "StockC 算定フロー（数値例：CLT 1 m³、耐用年数50年）")
    _fill_body(slide, [
        "算定手順:",
        ("① 建材 i のバイオジェニック炭素量 B_i [kgCO₂eq/m³] をFICES/INIESデータから取得", 1),
        ("② A1-A3（t=0）固定量:  −B_i × DCF(0)   = −B_i × 1.000", 1),
        ("③ C段階（t=t_EoL）放出: +B_i × DCF(t_EoL) = +B_i × 0.580（t_EoL=50年）", 1),
        ("④ 正味 StockC_i = B_i × [DCF(t_EoL) − 1] = B_i × (0.580 − 1.000) = −0.420 × B_i", 1),
        ("⑤ 建物全体: StockC = Σ StockC_i", 1),
        "",
        "数値例（CLT 1 m³、B_i = 269 kgCO₂eq）:",
    ], size=14)
    _add_table(slide,
        ["炭素フロー", "タイミング", "炭素量", "DCF係数", "補正後GHG"],
        [
            ["バイオジェニック炭素固定 A1-A3", "t=0",    "−269 kgCO₂eq", "1.000", "−269 kgCO₂eq"],
            ["解体焼却時の炭素放出 C4",         "t=50年", "+269 kgCO₂eq", "0.579", "+156 kgCO₂eq"],
            ["正味 StockC",                    "—",      "—",             "—",    "−113 kgCO₂eq/m³"],
        ],
        Emu(400_000), Emu(4_500_000), Emu(9_900_000), Emu(1_800_000),
        hsize=12, bsize=12
    )
    _add_chart(slide, "fig10_biogenic_carbon_cycle.png",
               Emu(6_500_000), Emu(1_900_000), Emu(3_200_000), Emu(2_400_000))
    _key_box(slide, "SFM前提（持続可能な森林管理＝FSC/PEFC認証）が成立することがStockCクレジットの前提条件")
    return slide


def s15_do100(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "評価期間 DO=100年 の選択根拠と批判")
    _add_table(slide,
        ["DO（観測期間）", "t=50年放出の係数", "木材への有利性", "備考"],
        [
            ["50年",             "0.000（期間外）",  "最大（放出が評価外）",  ""],
            ["100年（RE2020）",  "0.579",            "中程度",               "IPCCのGWP100と整合"],
            ["150年",            "0.721",            "低め",                 ""],
            ["∞（静的LCA相当）", "1.000",            "ゼロ",                 "動的補正なし"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(6_200_000), Emu(2_100_000),
        hsize=12, bsize=12
    )
    _fill_body(slide, [
        "RE2020が100年を選んだ理由: GWP100との整合性・Levasseur(2010)準拠・産業界への説明容易性",
        "",
        "Ventura & Feraille (2021) の批判:",
        ("t=50年の排出を0.58で評価するが、物理的に正確な値は0.72 → 約24%の過大評価", 1),
        ("t=75年超では〜53%の過大評価（0.580固定が原因）", 1),
        "",
        "ISO 21391-1:2025（フランス主導策定）: 変動DO方式（各排出フローに応じた可変DO）を採用。",
        "フランスは2028年フェーズ3でISO 21391-1方式への移行を検討中（Rivaton 勧告 No.5）。",
    ], size=13, ph_idx=1)
    _add_chart(slide, "fig09_time_horizon_comparison.png",
               Emu(6_400_000), Emu(1_700_000), Emu(3_400_000), Emu(3_000_000))
    return slide


# ── Section Ⅴ ──────────────────────────────────────────────────────────

def s16_sec_ecosystem(prs):
    slide = _add_slide(prs, L_SEC_NAV)
    _set_title(slide, "Ⅴ．算定エコシステム", size=28)
    _fill_body(slide, ["Elodie三層構造・ハイブリッド設計・Ic開示制度"], size=18)
    return slide


def s17_elodie(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "フランスの算定エコシステム：三層構造")
    _add_table(slide,
        ["層", "主要要素", "概要"],
        [
            ["第1層\nデータベース", "INIES / FDES",
             "中央データポータル（約5,500件）。建材メーカーが作成するフランス版EPD。\nバイオジェニック炭素データを含む（一部）。FDES不在時はデフォルト値（保守的）を使用。"],
            ["第2層\n計算ソフト", "Elodie（CSTB開発・無償）",
             "シェア80%超。INIESからFDESを直接インポートし、Ic計算・StockC算定を自動化。\nDHUP認定の申請書類（PDF）を生成。BIM連携（IFC）も一部対応。"],
            ["第3層\n申請・認証", "建築確認申請＋認証機関",
             "Elodie出力PDFを申請書類に添付（法的義務）。Cerqual・Bureau Veritasが審査。\n任意認証：BBCA（Ic≦seuil×0.70）、NF HQE等。"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(3_400_000),
        hsize=13, bsize=11
    )
    _fill_body(slide, [
        "他国比較: ドイツ（ÖKOBAUDAT）・スウェーデン・英国・日本（IDEA-BIM）は全て静的LCAのみ。",
        "フランスのみが動的LCA対応の公式計算ツールを義務制度として運用している。",
    ], size=13, ph_idx=1)
    return slide


def s18_hybrid(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "ハイブリッド設計：静的LCA×動的係数のモジュール別適用")
    _fill_body(slide, [
        "RE2020は「純粋な動的LCA」ではなく「静的GWP値 × 時間係数」のハイブリッド方式:",
        ("① FDES/EPD（静的GWP100値）を入力", 1),
        ("② モジュール発生時刻 t_m に対応するF_RE2020(t)を乗算", 1),
        ("③ StockCを加算 → Ic_construction を算出", 1),
    ], size=14)
    _add_table(slide,
        ["モジュール", "内容", "t_m", "F_RE2020(t)"],
        [
            ["A1-A3", "製品段階（製造・輸送）",      "0年",   "1.000"],
            ["B1-B5", "使用段階（保守・修繕・改修）", "25年",  "0.790"],
            ["C1-C4", "廃棄段階（解体・最終処分）",   "50年",  "0.580"],
            ["D（任意）", "再利用・エネルギー回収",   "50年超", "Ic除外"],
            ["StockC",   "バイオジェニック炭素クレジット", "—", "−0.42×B_i（正味）"],
        ],
        Emu(400_000), Emu(3_500_000), Emu(9_900_000), Emu(2_700_000),
        hsize=13, bsize=12
    )
    _key_box(slide,
        "ハイブリッド設計の合理性: 既存の静的EPDインフラを流用し、動的LCI整備コストを最小化しながら補正を実現")
    return slide


def s19_ic_disclosure(prs):
    slide = _add_slide(prs, L_CON_NAV)
    _set_title(slide, "Ic値の開示フローと利活用")
    _fill_body(slide, [
        "【開示フロー】",
        ("①設計者がElodieでIc試算 → ②seuil値との適合チェック", 1),
        ("③建築確認申請に計算書を添付（法的義務）→ ④認証機関審査 → ⑤引渡し時に最終値確定", 1),
        "",
        "【金融・ESG分野での活用（進行中）】",
        ("グリーンボンド適格性: EUタクソノミー技術基準の参照指標（2026年確定予定）", 1),
        ("住宅ローン優遇: BBCA認証建物への金利割引（Crédit Agricole等で実施）", 1),
        ("ESG報告: 不動産ポートフォリオのScope 3排出量計算にIc値を活用", 1),
        "",
        "【任意認証ラベル】",
        ("BBCA（低炭素建築）: Ic ≦ seuil×0.70  /  E+C-ラベル: Ic ≦ seuil×0.85", 1),
    ], size=14)
    _key_box(slide, "Ic値は建築確認書類→金融評価→EU規制対応まで一貫して参照される「炭素パスポート」として機能")
    return slide


# ── Section Ⅵ ──────────────────────────────────────────────────────────

def s20_sec_limits(prs):
    slide = _add_slide(prs, L_SEC_BLU)
    _set_title(slide, "Ⅵ．課題・現行制度の限界・日本への示唆", size=26)
    _fill_body(slide, ["Rivaton (2025) 評価と国際的位置づけ、日本の優先課題"], size=18)
    return slide


def s21_limits_rivaton(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "RE2020の構造的限界とRivaton (2025) 主要勧告")
    _add_table(slide,
        ["課題", "内容", "深刻度", "対処（2028年以降）"],
        [
            ["DO固定による過大評価",
             "t=50年の排出を0.58評価（物理的正確値0.72）→24%過大",
             "中（方向性は正しい）",
             "ISO 21391-1（変動DO）への移行（勧告No.5）"],
            ["FDESカバレッジ不足",
             "デフォルト値依存が慢性化。設計者間のIc値ばらつき大",
             "高",
             "2027年までに約10,000件に倍増（勧告No.1）"],
            ["Dモジュール除外",
             "再利用・リサイクルの炭素便益がIcに反映されずサーキュラー設計の誘因が弱い",
             "高",
             "部分的Ic組み込みを2028年から試験実施（勧告No.10）"],
            ["SFM前提の脆弱性",
             "非認証材・気候変動による森林劣化でクレジット根拠が崩れるリスク",
             "高（長期）",
             "FSC/PEFC認証強化（勧告No.7・8）"],
            ["AR4→AR6更新未対応",
             "a₀: 0.217→0.200、τ₁: 172.9→394.4年。係数値が若干変化",
             "低（方法論的正当性は維持）",
             "2028年フェーズ3でAR6パラメータ更新（勧告No.9）"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(4_200_000),
        hsize=12, bsize=10
    )
    _set_notes(slide, "Rivaton (2025) は12項目の勧告を提示。最優先はFDES拡充（No.1）とDO固定問題の解消（No.5）。")
    return slide


def s22_international(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "国際的位置づけ：ISO / EPBD / Level(s)")
    _add_table(slide,
        ["枠組み", "RE2020との関係", "ステータス"],
        [
            ["ISO 21391-1:2025",
             "動的バイオジェニック炭素評価の初の国際規格。フランス(AFNOR)主導で策定。\nRE2020の科学的アプローチを非線形DCFで精緻化した「国際昇華版」",
             "2025年成立\nJIS化検討が必要"],
            ["EPBD 2024（EU建物エネルギー指令改正）",
             "2030年以降の全EU新築建物にWLC（ライフサイクル炭素）開示義務。\nRE2020がEU全域義務化のモデルケースとして明示的に参照される",
             "2024年5月成立\nEU全加盟国で法制化作業中"],
            ["Level(s)（EU共通持続可能性指標）",
             "指標6（ライフサイクルGWP）がRE2020のIc指標に直接対応。\nEN 15978準拠の静的LCA指向だが、フランス動的LCA手法の組み込みを議論中",
             "改訂議論中\n動的LCA対応の方向"],
            ["EN 15804+A2（建材EPD規格）",
             "−1/+1アプローチ（正味ゼロ）が基本。RE2020のStockCとは哲学的相違。\nCERIB（コンクリート産業）はこの矛盾を「非整合性」として批判",
             "改訂交渉中\n動的評価の組み込みが焦点"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(4_000_000),
        hsize=12, bsize=10
    )
    _key_box(slide, "「先に制度を持つ国が国際規格を書く」— フランスは ISO 21391-1 策定でこの戦略を実現した")
    return slide


def s23_japan(prs):
    slide = _add_slide(prs, L_CON_BLU)
    _set_title(slide, "日本への示唆：優先課題と現状ギャップ")
    _add_table(slide,
        ["優先課題", "フランスの教訓", "日本の現状（2026年）", "推奨アクション"],
        [
            ["① 時間地平の合意\n（DO設定）",
             "DO=100年固定は後に24%過大評価問題を招いた。\n最初から変動DO（ISO 21391-1）採用が望ましい",
             "DO設定に関する合意なし",
             "ISO 21391-1:2025のJIS化着手。\n利害関係者間の合意形成が最初のステップ"],
            ["② データ基盤整備\n（時系列炭素フロー）",
             "FDESカバレッジ不足でデフォルト値依存が慢性化。\nデータ基盤は規制化に先行して整備すべき",
             "IDEAに時系列炭素フローデータなし",
             "IDEAへの排出タイミング情報追加。\n国産CLT・製材のFDES相当データを優先整備"],
            ["③ 実証プログラム先行\n（E+C-相当）",
             "10年・2,200棟の実証なしに規制化せず。\n証拠基盤なき規制化は反発・混乱を招く",
             "動的LCA実証プログラムなし",
             "CASBEE改訂とセットで300〜500棟規模の\nパイロットプログラムを5年計画で設計"],
            ["④ 計算ツール整備\n（Elodie相当）",
             "Elodie無償提供が導入障壁を大幅に低減。\nツールなしでは動的LCAは机上論",
             "動的LCA対応の国産ツールなし",
             "IDEA-BIM等に動的LCA機能追加。\n国産建材のStockC自動算定機能の開発"],
        ],
        Emu(400_000), Emu(1_700_000), Emu(9_900_000), Emu(4_800_000),
        hsize=12, bsize=10
    )
    _set_notes(slide, "日本では静的GWP100の同じ失敗が現在進行中であり、国産CLTが「排出材」と評価されている。フランスは10年の準備を省略しなかった。")
    return slide


def s24_end(prs):
    slide = _add_slide(prs, L_END_BLU)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "一般社団法人サステナブル経営推進機構（SuMPO）"
        run.font.size = Pt(20)
        run.font.bold = True
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            lines = [
                "SuMPO 動的LCA調査報告書",
                "調査ブロック①②統合版　2026年6月",
            ]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
    return slide


# ── Main ───────────────────────────────────────────────────────────────

def main():
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))

    # Remove template placeholder slides (keep structure)
    while len(prs.slides) > 0:
        sldIdLst = prs.slides._sldIdLst
        sldId    = sldIdLst[0]
        rId      = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)

    builders = [
        s01_cover,
        s02_agenda,
        s03_sec_why,
        s04_gwp100_failure,
        s05_4conditions,
        s06_sec_re2020,
        s07_3indicators,
        s08_seuil,
        s09_sec_theory,
        s10_dcf_bern,
        s11_0842,
        s12_sec_biogenic,
        s13_static_vs_dynamic,
        s14_stockc_calc,
        s15_do100,
        s16_sec_ecosystem,
        s17_elodie,
        s18_hybrid,
        s19_ic_disclosure,
        s20_sec_limits,
        s21_limits_rivaton,
        s22_international,
        s23_japan,
        s24_end,
    ]

    for fn in builders:
        fn(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"✅  Saved: {OUTPUT}")
    print(f"    Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
